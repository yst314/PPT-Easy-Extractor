import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image # Pillowライブラリの一部
import os
import threading # GUIがフリーズしないように処理を別スレッドで行う

class PptxExporterApp:
    def __init__(self, master):
        self.master = master
        master.title("PowerPoint自動出力システム")
        master.geometry("600x450") # ウィンドウサイズ調整

        # --- 変数定義 ---
        self.ppt_file_path = tk.StringVar()
        self.output_folder_path = tk.StringVar()

        # --- UI要素の作成 ---

        # PowerPointファイル選択
        tk.Label(master, text="PowerPointファイル:").grid(row=0, column=0, padx=5, pady=5, sticky="w")
        self.ppt_file_entry = tk.Entry(master, textvariable=self.ppt_file_path, width=60)
        self.ppt_file_entry.grid(row=0, column=1, padx=5, pady=5, sticky="ew")
        self.ppt_select_button = tk.Button(master, text="選択...", command=self.select_ppt_file)
        self.ppt_select_button.grid(row=0, column=2, padx=5, pady=5)

        # 出力先フォルダ選択
        tk.Label(master, text="出力先フォルダ:").grid(row=1, column=0, padx=5, pady=5, sticky="w")
        self.output_folder_entry = tk.Entry(master, textvariable=self.output_folder_path, width=60)
        self.output_folder_entry.grid(row=1, column=1, padx=5, pady=5, sticky="ew")
        self.output_folder_select_button = tk.Button(master, text="選択...", command=self.select_output_folder)
        self.output_folder_select_button.grid(row=1, column=2, padx=5, pady=5)

        # 実行ボタン
        self.run_button = tk.Button(master, text="実行", command=self.start_processing_thread, width=15, height=2)
        self.run_button.grid(row=2, column=1, padx=5, pady=20, sticky="s")

        # ステータス表示エリア
        tk.Label(master, text="ステータス:").grid(row=3, column=0, padx=5, pady=5, sticky="nw")
        self.status_text = scrolledtext.ScrolledText(master, wrap=tk.WORD, width=70, height=15)
        self.status_text.grid(row=4, column=0, columnspan=3, padx=5, pady=5, sticky="nsew")
        self.status_text.configure(state='disabled') # 初期状態は編集不可

        # グリッドの列の伸縮設定
        master.grid_columnconfigure(1, weight=1)
        # グリッドの行の伸縮設定（ステータスエリアを伸縮させるため）
        master.grid_rowconfigure(4, weight=1)


    def select_ppt_file(self):
        filepath = filedialog.askopenfilename(
            title="PowerPointファイルを選択",
            filetypes=(("PowerPointファイル", "*.pptx"), ("すべてのファイル", "*.*"))
        )
        if filepath:
            self.ppt_file_path.set(filepath)
            self.update_status(f"PowerPointファイル選択: {filepath}")

    def select_output_folder(self):
        folderpath = filedialog.askdirectory(title="出力先フォルダを選択")
        if folderpath:
            self.output_folder_path.set(folderpath)
            self.update_status(f"出力先フォルダ選択: {folderpath}")

    def update_status(self, message):
        self.status_text.configure(state='normal') # 編集可能に
        self.status_text.insert(tk.END, message + "\n")
        self.status_text.see(tk.END) # 自動スクロール
        self.status_text.configure(state='disabled') # 再び編集不可に
        self.master.update_idletasks() # GUIの更新を即時反映

    def start_processing_thread(self):
        """ 処理を別スレッドで実行する """
        ppt_path = self.ppt_file_path.get()
        output_path = self.output_folder_path.get()

        if not ppt_path:
            messagebox.showerror("エラー", "PowerPointファイルを選択してください。")
            return
        if not os.path.exists(ppt_path):
            messagebox.showerror("エラー", f"指定されたPowerPointファイルが見つかりません:\n{ppt_path}")
            return
        if not output_path:
            messagebox.showerror("エラー", "出力先フォルダを選択してください。")
            return

        # 出力先フォルダが存在しない場合は作成
        if not os.path.exists(output_path):
            try:
                os.makedirs(output_path)
                self.update_status(f"出力先フォルダを作成しました: {output_path}")
            except OSError as e:
                messagebox.showerror("エラー", f"出力先フォルダの作成に失敗しました:\n{output_path}\n{e}")
                return

        # 処理中はボタンを無効化
        self.run_button.config(state=tk.DISABLED)
        self.ppt_select_button.config(state=tk.DISABLED)
        self.output_folder_select_button.config(state=tk.DISABLED)

        # スレッドを作成して処理を開始
        thread = threading.Thread(target=self.process_presentation, args=(ppt_path, output_path))
        thread.daemon = True  # メインスレッド終了時に子スレッドも終了する
        thread.start()

    def process_presentation(self, ppt_file_path, output_dir):
        try:
            self.update_status("処理を開始します...")
            prs = Presentation(ppt_file_path)
            total_slides = len(prs.slides)
            self.update_status(f"合計スライド数: {total_slides}")

            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                self.update_status(f"--- スライド {slide_number}/{total_slides} の処理を開始 ---")

                # テキストの抽出と保存
                slide_text_content = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text_content.append(run.text)

                if slide_text_content:
                    text_file_name = f"スライド{slide_number}_text.txt"
                    text_file_path = os.path.join(output_dir, text_file_name)
                    try:
                        with open(text_file_path, "w", encoding="utf-8") as f:
                            f.write("\n".join(slide_text_content))
                        self.update_status(f"  テキストを保存: {text_file_name}")
                    except Exception as e:
                        self.update_status(f"  エラー: テキストファイルの保存に失敗しました - {text_file_name} ({e})")
                else:
                    self.update_status(f"  スライド {slide_number} には抽出可能なテキストがありませんでした。")


                # 画像の抽出と保存
                image_counter = 1
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext.lower() # 拡張子を小文字に

                            # Pillowがサポートする形式か確認 (基本的なもの)
                            supported_formats = ["jpg", "jpeg", "png", "gif", "bmp", "tiff"]
                            if image_ext not in supported_formats:
                                # サポート外の場合、PNGとして保存試行
                                self.update_status(f"  警告: 画像形式 '{image_ext}' は直接サポートされていません。PNGとして保存を試みます。")
                                image_ext = "png" # デフォルトでpng

                            image_filename = f"スライド{slide_number}_image{image_counter}.{image_ext}"
                            image_path = os.path.join(output_dir, image_filename)

                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            self.update_status(f"  画像を保存: {image_filename}")
                            image_counter += 1
                        except Exception as e:
                            self.update_status(f"  エラー: スライド {slide_number} の画像抽出/保存に失敗しました ({e})")
                    # TODO: グループ化された図形内の画像や、他の種類の埋め込みオブジェクト(OLEなど)内の画像抽出は
                    #       より複雑な処理が必要となり、python-pptxだけでは難しい場合がある。

                self.update_status(f"--- スライド {slide_number}/{total_slides} の処理を完了 ---")

            self.update_status("全ての処理が完了しました。")
            messagebox.showinfo("完了", "PowerPointファイルの処理が完了しました。")

        except FileNotFoundError:
            self.update_status(f"エラー: 指定されたPowerPointファイルが見つかりません - {ppt_file_path}")
            messagebox.showerror("エラー", f"指定されたPowerPointファイルが見つかりません:\n{ppt_file_path}")
        except Exception as e:
            self.update_status(f"エラーが発生しました: {e}")
            messagebox.showerror("エラー", f"処理中にエラーが発生しました:\n{e}")
        finally:
            # 処理が完了またはエラー発生後、ボタンを有効に戻す
            self.run_button.config(state=tk.NORMAL)
            self.ppt_select_button.config(state=tk.NORMAL)
            self.output_folder_select_button.config(state=tk.NORMAL)


if __name__ == "__main__":
    root = tk.Tk()
    app = PptxExporterApp(root)
    root.mainloop()